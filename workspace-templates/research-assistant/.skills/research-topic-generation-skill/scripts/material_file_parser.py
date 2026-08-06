#!/usr/bin/env python3
"""
研究选题生成 Skill - 材料文件解析模块

支持解析教师上传的四种文件类型：
  - PDF (.pdf)  — 提取全文文本，保留页码边界
  - PPTX (.pptx) — 提取幻灯片文本，保留幻灯片序号
  - DOCX (.docx) — 提取段落和表格文本
  - TXT  (.txt)  — 直接读取文本内容

使用方式：
  from material_file_parser import parse_file, batch_parse_files, ParsedFile

  result = parse_file("/path/to/file.pdf")
  print(result.text)           # 提取的全文
  print(result.metadata)       # 文件元信息（页数、段落数等）

  results = batch_parse_files(["/path/a.pdf", "/path/b.docx"])
  for r in results:
      print(r.filename, r.text[:100])
"""

from __future__ import annotations

import os
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 数据模型
# ---------------------------------------------------------------------------

@dataclass
class ParsedFile:
    """单个文件的解析结果。"""
    filepath: str
    filename: str
    file_ext: str                          # 扩展名（含点），如 ".pdf"
    text: str                              # 提取的全文文本
    metadata: dict[str, Any] = field(default_factory=dict)
    # metadata 示例：
    #   PDF:  {"pages": 5, "title": "..."}
    #   PPTX: {"slides": 12, "slide_texts": [...]}
    #   DOCX: {"paragraphs": 30, "tables": 2}
    #   TXT:  {"encoding": "utf-8"}
    errors: list[str] = field(default_factory=list)

    @property
    def text_length(self) -> int:
        return len(self.text)

    @property
    def is_empty(self) -> bool:
        return not self.text.strip()

    def summary(self) -> str:
        """返回人类可读的解析摘要。"""
        meta_str = ", ".join(f"{k}={v}" for k, v in self.metadata.items())
        preview = self.text[:120].replace("\n", " ").strip()
        return (
            f"[{self.filename}] "
            f"类型={self.file_ext}, "
            f"文本长度={self.text_length} 字符, "
            f"元信息=({meta_str}), "
            f"错误={len(self.errors)}, "
            f"预览: {preview}..."
        )


# ---------------------------------------------------------------------------
# 格式解析器
# ---------------------------------------------------------------------------

def _parse_pdf(filepath: str) -> ParsedFile:
    """使用 PyPDF2 解析 PDF 文件。"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return _fallback_binary(filepath, ".pdf",
                                "PyPDF2 未安装。请运行: pip install PyPDF2")

    metadata: dict[str, Any] = {}
    errors: list[str] = []
    pages_text: list[str] = []

    try:
        reader = PdfReader(filepath)
        num_pages = len(reader.pages)
        metadata["pages"] = num_pages

        # 尝试读文档信息
        if reader.metadata:
            doc_info = reader.metadata
            if doc_info.title:
                metadata["title"] = str(doc_info.title)
            if doc_info.author:
                metadata["author"] = str(doc_info.author)

        for i, page in enumerate(reader.pages, 1):
            try:
                page_text = page.extract_text()
                if page_text:
                    pages_text.append(page_text.strip())
            except Exception as e:
                errors.append(f"第{i}页解析失败: {e}")

        full_text = "\n\n".join(pages_text)
        metadata["pages_with_text"] = len(pages_text)

    except Exception as e:
        errors.append(f"PDF 文件解析失败: {e}")
        full_text = ""

    return ParsedFile(
        filepath=filepath,
        filename=os.path.basename(filepath),
        file_ext=".pdf",
        text=full_text,
        metadata=metadata,
        errors=errors,
    )


def _parse_pptx(filepath: str) -> ParsedFile:
    """使用 python-pptx 解析 PPTX 文件。"""
    try:
        from pptx import Presentation
    except ImportError:
        return _fallback_binary(filepath, ".pptx",
                                "python-pptx 未安装。请运行: pip install python-pptx")

    metadata: dict[str, Any] = {}
    errors: list[str] = []
    slide_texts: list[str] = []
    all_shapes_count = 0

    try:
        prs = Presentation(filepath)
        num_slides = len(prs.slides)
        metadata["slides"] = num_slides

        for i, slide in enumerate(prs.slides, 1):
            shapes_text: list[str] = []
            for shape in slide.shapes:
                all_shapes_count += 1
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            shapes_text.append(para_text)
                if shape.has_table:
                    table_text = _extract_table_text(shape.table)
                    if table_text:
                        shapes_text.append(table_text)
            slide_text = "\n".join(shapes_text)
            if slide_text:
                slide_texts.append(f"[幻灯片 {i}]\n{slide_text}")

        full_text = "\n\n".join(slide_texts)
        metadata["slide_texts"] = slide_texts
        metadata["shapes_total"] = all_shapes_count
        metadata["slides_with_text"] = len(slide_texts)

    except Exception as e:
        errors.append(f"PPTX 文件解析失败: {e}")
        full_text = ""

    return ParsedFile(
        filepath=filepath,
        filename=os.path.basename(filepath),
        file_ext=".pptx",
        text=full_text,
        metadata=metadata,
        errors=errors,
    )


def _parse_docx(filepath: str) -> ParsedFile:
    """使用 python-docx 解析 DOCX 文件。"""
    try:
        from docx import Document
    except ImportError:
        return _fallback_binary(filepath, ".docx",
                                "python-docx 未安装。请运行: pip install python-docx")

    metadata: dict[str, Any] = {}
    errors: list[str] = []
    paragraphs_text: list[str] = []
    tables_text: list[str] = []

    try:
        doc = Document(filepath)

        # 段落文本
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs_text.append(text)

        # 表格文本
        for table in doc.tables:
            table_content = _extract_table_text(table)
            if table_content:
                tables_text.append(table_content)

        metadata["paragraphs"] = len(doc.paragraphs)
        metadata["paragraphs_with_text"] = len(paragraphs_text)
        metadata["tables"] = len(doc.tables)
        metadata["tables_with_text"] = len(tables_text)

        parts = paragraphs_text + tables_text
        full_text = "\n\n".join(parts)

    except Exception as e:
        errors.append(f"DOCX 文件解析失败: {e}")
        full_text = ""

    return ParsedFile(
        filepath=filepath,
        filename=os.path.basename(filepath),
        file_ext=".docx",
        text=full_text,
        metadata=metadata,
        errors=errors,
    )


def _parse_txt(filepath: str) -> ParsedFile:
    """解析纯文本文件。"""
    metadata: dict[str, Any] = {}
    errors: list[str] = []
    full_text = ""

    for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            with open(filepath, "r", encoding=encoding) as f:
                full_text = f.read()
            metadata["encoding"] = encoding
            break
        except UnicodeDecodeError:
            continue
        except Exception as e:
            errors.append(f"读取失败 ({encoding}): {e}")

    if not full_text and not errors:
        errors.append("无法以任何编码读取文件。")

    metadata["lines"] = full_text.count("\n") + 1 if full_text else 0

    return ParsedFile(
        filepath=filepath,
        filename=os.path.basename(filepath),
        file_ext=".txt",
        text=full_text,
        metadata=metadata,
        errors=errors,
    )


def _fallback_binary(filepath: str, ext: str, reason: str) -> ParsedFile:
    """缺少解析库时的降级结果。"""
    return ParsedFile(
        filepath=filepath,
        filename=os.path.basename(filepath),
        file_ext=ext,
        text="",
        metadata={},
        errors=[reason],
    )


def _extract_table_text(table: Any) -> str:
    """从 DOCX/PPTX 表格中提取文本（行列格式）。"""
    rows: list[str] = []
    try:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
    except Exception:
        pass
    return "\n".join(rows)


# ---------------------------------------------------------------------------
# 解析器调度表
# ---------------------------------------------------------------------------

_PARSER_MAP: dict[str, Any] = {
    ".pdf": _parse_pdf,
    ".pptx": _parse_pptx,
    ".docx": _parse_docx,
    ".txt": _parse_txt,
}

_SUPPORTED_EXTS = frozenset(_PARSER_MAP.keys())


def supported_extensions() -> list[str]:
    """返回支持的文件扩展名列表。"""
    return sorted(_SUPPORTED_EXTS)


def is_supported(filepath: str | Path) -> bool:
    """判断文件扩展名是否被支持。"""
    ext = os.path.splitext(str(filepath))[1].lower()
    return ext in _SUPPORTED_EXTS


# ---------------------------------------------------------------------------
# 公开 API
# ---------------------------------------------------------------------------

def parse_file(filepath: str | Path) -> ParsedFile:
    """解析单个文件，自动识别格式。

    Args:
        filepath: 文件路径

    Returns:
        ParsedFile 对象，包含提取的文本和元信息。
        如果格式不支持或解析失败，text 为空字符串，errors 中记录原因。

    Raises:
        FileNotFoundError: 文件不存在
    """
    filepath_str = str(filepath)

    if not os.path.isfile(filepath_str):
        raise FileNotFoundError(f"文件不存在: {filepath_str}")

    ext = os.path.splitext(filepath_str)[1].lower()
    parser = _PARSER_MAP.get(ext)

    if parser is None:
        return ParsedFile(
            filepath=filepath_str,
            filename=os.path.basename(filepath_str),
            file_ext=ext,
            text="",
            metadata={},
            errors=[f"不支持的文件格式: {ext}。支持的格式: {sorted(_SUPPORTED_EXTS)}"],
        )

    logger.info("解析文件: %s (格式=%s)", filepath_str, ext)
    result = parser(filepath_str)
    logger.info("解析完成: %s", result.summary())
    return result


def batch_parse_files(filepaths: list[str | Path]) -> list[ParsedFile]:
    """批量解析多个文件。

    Args:
        filepaths: 文件路径列表

    Returns:
        ParsedFile 列表，顺序与输入一致。
    """
    results: list[ParsedFile] = []
    for fp in filepaths:
        try:
            results.append(parse_file(fp))
        except FileNotFoundError as e:
            results.append(ParsedFile(
                filepath=str(fp),
                filename=os.path.basename(str(fp)),
                file_ext=os.path.splitext(str(fp))[1].lower(),
                text="",
                metadata={},
                errors=[str(e)],
            ))
    return results


def parsed_to_source_material(
    parsed: ParsedFile,
    material_id: str | None = None,
    material_type: str | None = None,
) -> dict[str, Any]:
    """将 ParsedFile 转换为 SourceMaterial 字典（供 material_adapter 使用）。

    Args:
        parsed: 解析结果
        material_id: 材料 ID，不提供则用文件名
        material_type: 材料类型，不提供则根据扩展名推断

    Returns:
        符合 source_material.schema.json 的字典
    """
    ext = parsed.file_ext.lower()

    # 文件扩展名 → 材料类型映射
    type_map = {
        ".pdf": "uploaded_paper_text",
        ".docx": "lesson_case",
        ".pptx": "other",           # PPT 通常是课件/分享
        ".txt": "teaching_reflection",
    }
    default_type = type_map.get(ext, "other")

    return {
        "materialId": material_id or f"file-{parsed.filename}",
        "materialType": material_type or default_type,
        "title": os.path.splitext(parsed.filename)[0],
        "content": parsed.text,
        "rawText": parsed.text,
        "filePath": parsed.filepath,
        "sourceStatus": "user_provided",
        "sensitivity": "internal",
        "permissions": {
            "canUseForGeneration": True,
            "canStore": True,
            "canExport": True,
            "limits": ["用户材料事实需由教师确认。"],
        },
        "metadata": {
            "fileType": ext,
            "parserMetadata": parsed.metadata,
            "parseErrors": parsed.errors,
        },
    }


def batch_to_source_materials(
    parsed_files: list[ParsedFile],
    material_type_map: dict[str, str] | None = None,
) -> list[dict[str, Any]]:
    """批量将 ParsedFile 转换为 SourceMaterial 列表。

    Args:
        parsed_files: 解析结果列表
        material_type_map: 文件名 → 材料类型的映射（可选，用于覆盖默认推断）

    Returns:
        SourceMaterial 字典列表
    """
    materials: list[dict[str, Any]] = []
    for i, parsed in enumerate(parsed_files, 1):
        override_type = None
        if material_type_map:
            override_type = material_type_map.get(parsed.filename)
        materials.append(
            parsed_to_source_material(
                parsed,
                material_id=f"mat-{i:03d}",
                material_type=override_type,
            )
        )
    return materials


# ---------------------------------------------------------------------------
# CLI 调试入口
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print(f"用法: python {sys.argv[0]} <file1> [file2 ...]")
        print(f"支持格式: {sorted(_SUPPORTED_EXTS)}")
        sys.exit(1)

    filepaths = sys.argv[1:]
    results = batch_parse_files(filepaths)

    for r in results:
        print(r.summary())
        if r.errors:
            for err in r.errors:
                print(f"  [错误] {err}")
        print("-" * 60)
